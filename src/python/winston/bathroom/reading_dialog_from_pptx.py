# 02-08-2018
# based on http://python-pptx.readthedocs.io/en/latest/user/quickstart.html
# for regulare expressions
# https://www.datacamp.com/community/tutorials/python-regular-expression-tutorial?utm_source=adwords_ppc&utm_campaignid=898687156&utm_adgroupid=48947256715&utm_device=c&utm_keyword=&utm_matchtype=b&utm_network=g&utm_adpostion=1t1&utm_creative=255798340456&utm_targetid=aud-299261629574:dsa-473406579035&utm_loc_interest_ms=&utm_loc_physical_ms=1010754&gclid=CjwKCAjw14rbBRB3EiwAKeoG_67in5hOsRcUn49QWBMj5qybxE3S-Mo1uJj5GVHGamJ8KaMZ3L66_xoColIQAvD_BwE

from pptx import Presentation
import re
import numpy as np

def main():
    print("hallo")

    prs = Presentation("Bathroom_prr.pptx")

    # text_runs wil be populated with a list of strings,
    # one for each text run in the presentation
    text_runs = []
    dialog_runs = []

    bullet_slide_layout = prs.slide_layouts[1]

    slide_nr = 1

    pattern = re.compile(r'^Winston:|^Customer:', re.IGNORECASE)

    for slide in prs.slides:
        print("slide %d" % (slide_nr))

        for shape in slide.placeholders:
            print('%d %s %s' % (shape.placeholder_format.idx, shape.name, shape.shape_type))
            if shape.is_placeholder:
                phf = shape.placeholder_format
                print('\t%d, %s' % (phf.idx, phf.type))
                # look for bullets and get all the text from this section
                if phf.type == 7:
                    print("\tyes")
                    if not shape.has_text_frame:
                        continue
                    line = ""
                    for paragraph in shape.text_frame.paragraphs:
                        # if not line:
                        #     print("empty")
                        for run in paragraph.runs:
                            #raw_text = "%r"%run.text
                            #print("\t\t%r" % (run.text))
                            #match = re.search(r'^Winston:|^Customer:', run.text)
                            if pattern.match(run.text):
                                # match so new dialog sentence start , store the last one if any
                                if not line:
                                    print("empty")
                                else:
                                    dialog_runs.append(line)
                                # start building new dialog
                                line = run.text
                            else:
                                # not a winston or customer so concat
                                line = line + run.text





            # if not shape.has_text_frame:
            #     continue
            # for paragraph in shape.text_frame.paragraphs:
            #     for run in paragraph.runs:
            #         # convert in raw text
            #         line = run.text
            #         #print(line)
            #         if re.match(r'^Winston:', line) or re.match(r'Customer:', line):
            #             dialog_runs.append(line)
            #         text_runs.append(line)
        slide_nr = slide_nr + 1
        print(" ")

    #print(text_runs)
    np.savetxt("text.txt", text_runs, delimiter=" ", newline="\n", fmt="%s")
    np.savetxt("dialog.txt", dialog_runs, delimiter=" ", newline="\n", fmt="%s")



if __name__ == "__main__":
    main()